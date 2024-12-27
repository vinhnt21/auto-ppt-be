from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE
from pptx.dml.color import RGBColor
import markdown2
from bs4 import BeautifulSoup
import requests
from io import BytesIO
import re


def hex_to_rgb(hex_color):
    hex_color = hex_color.lstrip("#")
    return tuple(int(hex_color[i : i + 2], 16) for i in (0, 2, 4))


def setup_text_frame(text_frame):
    """Helper function to setup text frame properties"""
    text_frame.word_wrap = True
    text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    return text_frame


def markdown_to_pptx(markdown_content, output_file, custom_theme=None):
    # Default theme
    default_theme = {
        "FONT": "Arial",
        "TITLE_COLOR": "#083B4A",
        "CONTENT_COLOR": "#000000",
        "BG_COLOR": "#D3D3D3",
    }

    theme = default_theme if custom_theme is None else {**default_theme, **custom_theme}

    # Create presentation
    prs = Presentation()
    prs.slide_width = Inches(16)
    prs.slide_height = Inches(9)

    # Convert markdown to HTML
    html_content = markdown2.markdown(markdown_content)
    soup = BeautifulSoup(html_content, "html.parser")

    # Split content into slides
    slides = []
    current_slide = []

    for element in soup.find_all(["h2", "ul", "img"]):
        if element.name == "h2" and current_slide:
            slides.append(current_slide)
            current_slide = []
        current_slide.append(element)
    if current_slide:
        slides.append(current_slide)

    # Process slides
    for i, slide_content in enumerate(slides):
        if i == 0:  # First slide
            slide = prs.slides.add_slide(prs.slide_layouts[5])

            # Add title
            title_box = slide.shapes.add_textbox(
                Inches(1), Inches(2), Inches(14), Inches(2)
            )
            title_frame = setup_text_frame(title_box.text_frame)
            title_para = title_frame.add_paragraph()
            title_para.text = slide_content[0].text
            title_para.font.size = Pt(36)
            title_para.font.name = theme["FONT"]
            title_para.font.color.rgb = RGBColor(*hex_to_rgb(theme["TITLE_COLOR"]))

            # Add subtitle if exists
            if len(slide_content) > 1 and slide_content[1].name == "ul":
                subtitle_box = slide.shapes.add_textbox(
                    Inches(1), Inches(4), Inches(14), Inches(1)
                )
                subtitle_frame = setup_text_frame(subtitle_box.text_frame)
                subtitle_para = subtitle_frame.add_paragraph()
                subtitle_para.text = slide_content[1].find("li").text
                subtitle_para.font.size = Pt(24)
                subtitle_para.font.name = theme["FONT"]
                subtitle_para.font.color.rgb = RGBColor(
                    *hex_to_rgb(theme["CONTENT_COLOR"])
                )
                subtitle_box.text_frame.word_wrap = True
                
        elif i == len(slides) - 1:  # Last slide
            slide = prs.slides.add_slide(prs.slide_layouts[5])

            # Add centered title
            title_box = slide.shapes.add_textbox(
                Inches(1), Inches(3.5), Inches(14), Inches(2)
            )
            title_frame = setup_text_frame(title_box.text_frame)
            title_para = title_frame.add_paragraph()
            title_para.text = slide_content[0].text
            title_para.alignment = PP_ALIGN.CENTER
            title_para.font.size = Pt(36)
            title_para.font.name = theme["FONT"]
            title_para.font.color.rgb = RGBColor(*hex_to_rgb(theme["TITLE_COLOR"]))

        else:  # Content slides
            slide = prs.slides.add_slide(prs.slide_layouts[5])

            # Add title
            title_box = slide.shapes.add_textbox(
                Inches(1), Inches(0.5), Inches(14), Inches(1)
            )
            title_frame = setup_text_frame(title_box.text_frame)
            title_para = title_frame.add_paragraph()
            title_para.text = slide_content[0].text
            title_para.font.size = Pt(36)
            title_para.font.name = theme["FONT"]
            title_para.font.color.rgb = RGBColor(*hex_to_rgb(theme["TITLE_COLOR"]))

            # Determine content and image positions based on slide number
            if i % 2 == 0:  # Even slides: content left, image right
                content_left, content_width = Inches(1), Inches(7)
                image_left, image_width = Inches(8.5), Inches(6)
            else:  # Odd slides: image left, content right
                image_left, image_width = Inches(1.5), Inches(6)
                content_left, content_width = Inches(8), Inches(7)

            # Add content
            content_box = slide.shapes.add_textbox(
                content_left, Inches(2), content_width, Inches(6)
            )
            content_frame = setup_text_frame(content_box.text_frame)

            for element in slide_content[1:]:
                if element.name == "ul":
                    for li in element.find_all("li"):
                        p = content_frame.add_paragraph()
                        p.text = "• " + li.text
                        p.font.size = Pt(24)
                        p.font.name = theme["FONT"]
                        p.font.color.rgb = RGBColor(*hex_to_rgb(theme["CONTENT_COLOR"]))
                        p.line_spacing = (
                            1.2  # Add some line spacing for better readability
                        )
                elif element.name == "img":
                    img_src = element.get("src")
                    if img_src:
                        try:
                            if img_src.startswith("http"):
                                response = requests.get(img_src)
                                image_stream = BytesIO(response.content)
                            else:
                                image_stream = img_src
                            slide.shapes.add_picture(
                                image_stream, image_left, Inches(2), image_width
                            )
                        except Exception as e:
                            print(f"Error adding image: {e}")

    # Save presentation
    prs.save(output_file)
    print(f"Created PowerPoint file: {output_file}")
