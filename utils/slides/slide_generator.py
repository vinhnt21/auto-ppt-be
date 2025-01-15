from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_AUTO_SIZE
from pptx.util import Pt
from pptx.util import Inches
import requests
from io import BytesIO
from env_config import SLIDE_FOLDER

# Layout mapping
LAYOUT = {
    "Title Slide": 0,
    "Title and Content": 1,
    "Section Header": 2,
    "Two Content": 3,
    "Comparison": 4,
    "Content with Caption": 7,
    "Picture with Caption": 8,
}


def get_layout(name: str) -> int:
    name = name.strip()
    if name not in LAYOUT:
        raise ValueError(f"Invalid slide layout: {name}")
    return LAYOUT[name]


def init_presentation(template_name=None):
    """
    Khởi tạo đối tượng Presentation với tỷ lệ khung hình 16:9.
    Args:
        template_path: Đường dẫn tới template (mặc định: None)
    Returns:
        prs: Đối tượng Presentation
    """
    THEME_TEMPLATES = {
        "Office": "templates/office_theme.pptx",
        "Facet": "templates/facet_theme.pptx",
        "Integral": "templates/integral_theme.pptx",
        "Ion": "templates/ion_theme.pptx",
        "Vapor": "templates/vapor_theme.pptx",
        "Damask": "templates/damask_theme.pptx",
    }
    if template_name is None or template_name not in THEME_TEMPLATES:
        template_path = THEME_TEMPLATES["Facet"]
    template_path = THEME_TEMPLATES[template_name]
    prs = Presentation(template_path)
    # Đặt tỷ lệ khung hình 16:9 (10 x 5.625 inch)
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)

    return prs


# Các hàm tạo slide
def add_title_slide(prs, title: str, subtitle: str = ""):
    slide = prs.slides.add_slide(prs.slide_layouts[get_layout("Title Slide")])
    slide.shapes.title.text = title
    if subtitle:
        slide.placeholders[1].text = subtitle
    return slide


def add_title_and_content_slide(prs, title: str, content: str):
    slide = prs.slides.add_slide(prs.slide_layouts[get_layout("Title and Content")])
    slide.shapes.title.text = title
    slide.placeholders[1].text = content
    return slide


def add_section_header_slide(prs, title: str, subtitle: str = ""):
    slide = prs.slides.add_slide(prs.slide_layouts[get_layout("Section Header")])
    slide.shapes.title.text = title
    if subtitle:
        slide.placeholders[1].text = subtitle
    return slide


def add_two_content_slide(prs, title: str, left_content: str, right_content: str):
    slide = prs.slides.add_slide(prs.slide_layouts[get_layout("Two Content")])
    slide.shapes.title.text = title
    slide.placeholders[1].text = left_content
    slide.placeholders[2].text = right_content
    return slide


def add_comparison_slide(
    prs,
    title: str,
    left_title: str,
    left_content: str,
    right_title: str,
    right_content: str,
):
    slide = prs.slides.add_slide(prs.slide_layouts[get_layout("Comparison")])
    slide.shapes.title.text = title
    slide.placeholders[1].text = left_title
    slide.placeholders[2].text = left_content
    slide.placeholders[3].text = right_title
    slide.placeholders[4].text = right_content
    return slide


def add_content_with_caption_slide(prs, title: str, content: str, caption: str):
    slide = prs.slides.add_slide(prs.slide_layouts[get_layout("Content with Caption")])
    slide.shapes.title.text = title
    slide.placeholders[1].text = content
    slide.placeholders[2].text = caption
    return slide


def add_picture_with_caption_slide(prs, title: str, image_path: str, caption: str):
    """
    Thêm slide với hình ảnh và chú thích, hỗ trợ cả đường dẫn local và URL
    """

    def download_image(url):
        try:
            response = requests.get(url, timeout=10)
            response.raise_for_status()
            return BytesIO(response.content)
        except Exception as e:
            print(f"Error downloading image from {url}: {e}")
            return None

    slide = prs.slides.add_slide(prs.slide_layouts[get_layout("Picture with Caption")])
    slide.shapes.title.text = title
    placeholder = slide.placeholders[1]

    # Kiểm tra nếu image_path là URL
    if image_path.startswith(("http://", "https://")):
        image_data = download_image(image_path)
        if image_data:
            placeholder.insert_picture(image_data)
        else:
            # Có thể thêm hình ảnh placeholder hoặc thông báo lỗi
            placeholder.text = "Failed to load image"
    else:
        # Xử lý đường dẫn local như bình thường
        placeholder.insert_picture(image_path)

    slide.placeholders[2].text = caption
    return slide


def process_content(content: str) -> str:
    """
    Xử lý nội dung của slide.
    """
    content = content.strip()
    content = content.split("- ")
    content = [i.strip() for i in content if i]
    content = "\n".join(content)
    return content


def create_slide(prs, layout_name, **kwargs):
    """
    Hàm điều phối để tạo slide dựa trên layout.
    """
    layout_name = layout_name.strip()

    if kwargs.get("content"):
        kwargs["content"] = process_content(kwargs["content"])
    if kwargs.get("subtitle"):
        kwargs["subtitle"] = process_content(kwargs["subtitle"])
    if kwargs.get("left_content"):
        kwargs["left_content"] = process_content(kwargs["left_content"])
    if kwargs.get("right_content"):
        kwargs["right_content"] = process_content(kwargs["right_content"])

    if layout_name == "Title Slide":
        return add_title_slide(prs, kwargs.get("title", ""), kwargs.get("subtitle", ""))
    elif layout_name == "Title and Content":
        return add_title_and_content_slide(
            prs, kwargs.get("title", ""), kwargs.get("content", "")
        )
    elif layout_name == "Section Header":
        return add_section_header_slide(
            prs, kwargs.get("title", ""), kwargs.get("subtitle", "")
        )
    elif layout_name == "Two Content":
        return add_two_content_slide(
            prs,
            kwargs.get("title", ""),
            kwargs.get("left_content", ""),
            kwargs.get("right_content", ""),
        )
    elif layout_name == "Comparison":
        return add_comparison_slide(
            prs,
            kwargs.get("title", ""),
            kwargs.get("left_title", ""),
            kwargs.get("left_content", ""),
            kwargs.get("right_title", ""),
            kwargs.get("right_content", ""),
        )
    elif layout_name == "Content with Caption":
        return add_content_with_caption_slide(
            prs,
            kwargs.get("title", ""),
            kwargs.get("content", ""),
            kwargs.get("caption", ""),
        )
    elif layout_name == "Picture with Caption":
        return add_picture_with_caption_slide(
            prs,
            kwargs.get("title", ""),
            kwargs.get("image_path", ""),
            kwargs.get("caption", ""),
        )
    else:
        raise ValueError(f"Invalid layout name: {layout_name}")


def delete_slides_range(prs, start_index, end_index):
    """
    Xóa một dãy slides từ start_index đến end_index.

    Args:
        prs: Đối tượng Presentation
        start_index: Index bắt đầu (inclusive)
        end_index: Index kết thúc (inclusive)
    """

    def delete_slide_by_index(prs, index):
        """
        Xóa một slide dựa trên index.

        Args:
            prs: Đối tượng Presentation
            index: Index của slide cần xóa (0-based)
        """
        try:
            if not 0 <= index < len(prs.slides):
                raise ValueError(
                    f"Invalid slide index: {index}. Index must be between 0 and {len(prs.slides)-1}"
                )

            xml_slides = prs.slides._sldIdLst
            slides = list(xml_slides)
            xml_slides.remove(slides[index])
            return True
        except Exception as e:
            print(f"Error deleting slide at index {index}: {str(e)}")
            return False

    try:
        if not (
            0 <= start_index < len(prs.slides) and 0 <= end_index < len(prs.slides)
        ):
            raise ValueError(
                f"Invalid index range. Must be between 0 and {len(prs.slides)-1}"
            )
        if start_index > end_index:
            raise ValueError("Start index must be less than or equal to end index")

        # Xóa từ cuối lên để tránh index bị thay đổi
        for i in range(end_index, start_index - 1, -1):
            delete_slide_by_index(prs, i)
        return True
    except Exception as e:
        print(f"Error deleting slides range: {str(e)}")
        return False


def save_presentation(prs, output_path: str):
    """
    Lưu presentation.
    """
    LENGTH_OF_TEMPLATE = 9
    try:
        delete_slides_range(prs, 0, LENGTH_OF_TEMPLATE - 1)
        prs.save(f"{SLIDE_FOLDER}/{output_path}")
    except Exception as e:
        raise IOError(f"Failed to save presentation: {str(e)}")
